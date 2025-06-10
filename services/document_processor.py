import os
from typing import List, Dict, Any, Optional
import asyncio
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_huggingface import HuggingFaceEmbeddings

from logger import logger
from config import EMBEDDING_MODEL

# 임베딩 모델 초기화
embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

async def extract_text(file_path: str, file_ext: str) -> str:
    """파일에서 텍스트를 추출합니다."""
    
    # 파일 형식에 따라 적절한 추출 방법 사용
    if file_ext in ['.txt', '.md', '.csv']:
        # 텍스트 파일
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    elif file_ext in ['.pdf']:
        # PDF 파일
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n\n"
        return text
    
    elif file_ext in ['.docx', '.doc']:
        # Word 문서
        import docx
        
        doc = docx.Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    
    elif file_ext in ['.pptx', '.ppt']:
        # PowerPoint 파일
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    elif file_ext in ['.xlsx', '.xls']:
        # Excel 파일
        import pandas as pd
        
        # 모든 시트의 데이터를 텍스트로 변환
        dfs = pd.read_excel(file_path, sheet_name=None)
        text = ""
        for sheet_name, df in dfs.items():
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string(index=False) + "\n\n"
        return text
    
    else:
        # 지원되지 않는 형식
        logger.error(f"지원되지 않는 파일 형식: {file_ext}")
        return ""

async def process_document(
    text: str, 
    metadata: Dict[str, Any], 
    chunk_size: int = 500, 
    chunk_overlap: int = 50,
    embedding_model: str = None
) -> List[Document]:
    """문서를 청크로 나누고 임베딩합니다."""
    
    # 텍스트 분할기
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    
    # 텍스트를 청크로 분할
    chunks = text_splitter.create_documents(
        texts=[text],
        metadatas=[metadata]
    )
    
    # 각 청크에 인덱스 추가
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = i
    
    # 임베딩 생성 (비동기로 처리)
    async def embed_chunk(chunk):
        try:
            # 임베딩 생성 - HuggingFace 임베딩은 상대적으로 무거울 수 있으므로 비동기 처리
            embedding = await asyncio.to_thread(
                embeddings.embed_query,
                chunk.page_content
            )
            # 임베딩 저장
            chunk.embedding = embedding
            return chunk
        except Exception as e:
            logger.error(f"임베딩 생성 중 오류: {e}")
            # 오류 발생 시 빈 임베딩으로 처리
            chunk.embedding = [0.0] * embeddings.embedding_size
            return chunk
    
    # 모든 청크 임베딩을 병렬로 처리
    tasks = [embed_chunk(chunk) for chunk in chunks]
    embedded_chunks = await asyncio.gather(*tasks)
    
    return embedded_chunks